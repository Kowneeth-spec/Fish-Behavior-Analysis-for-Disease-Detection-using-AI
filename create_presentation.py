#!/usr/bin/env python3
"""
Fish Behavior Analysis for Disease Detection - PowerPoint Generator
Automatically creates a professional 21-slide presentation from PRESENTATION_PROMPT.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# Color definitions
BLUE = RGBColor(0, 105, 148)          # Ocean Blue
ORANGE = RGBColor(255, 140, 66)       # Orange
GREEN = RGBColor(46, 204, 113)        # Sea Green
DARK_GRAY = RGBColor(51, 51, 51)      # Dark Gray
LIGHT_GRAY = RGBColor(245, 245, 245)  # Light Gray

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle, notes=""):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(60)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_frame.text = subtitle
    for p in subtitle_frame.paragraphs:
        p.font.size = Pt(32)
        p.font.color.rgb = ORANGE
        p.alignment = PP_ALIGN.CENTER
    
    # Notes
    if notes:
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = notes
    
    return slide

def add_content_slide(title, content_list, notes=""):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_GRAY
    
    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = BLUE
    title_shape.line.color.rgb = BLUE
    
    # Title text
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.space_before = Pt(12)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(8.5), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.font.size = Pt(24)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)
        p.level = 0
    
    # Notes
    if notes:
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = notes
    
    return slide

# ============================================================
# SECTION 1: TITLE & INTRODUCTION (Slides 1-3)
# ============================================================

add_title_slide(
    "🐟 Fish Behavior Analysis for Disease Detection Using AI",
    "Enhancing Fish Welfare in Aquaculture Through Computer Vision",
    """Welcome to Fish Behavior Analysis for Disease Detection.
    
Today we'll explore an innovative project combining AI and aquaculture to revolutionize 
fish health monitoring. We'll discuss the problem, how the technology works, real-world 
implementation, and exciting opportunities ahead."""
)

add_content_slide(
    "Challenge: Fish Disease Detection in Aquaculture",
    [
        "❌ Traditional Methods Are Inadequate",
        "   • Manual observation: Time-consuming, inconsistent",
        "   • Blood tests: Invasive, expensive, late detection",
        "   • Economic impact: $10B+ annual losses globally",
        "",
        "❌ Disease Spread is Rapid",
        "   • Farm density: 100s of fish in small spaces",
        "   • Infection timeline: 24-48 hours to population spread",
        "   • Manual detection: Too slow to prevent outbreak",
        "",
        "❌ Welfare Concerns",
        "   • Stress indicators often invisible",
        "   • Early behavioral changes = disease warning signs"
    ],
    """Traditional methods for detecting fish disease are inadequate. Manual observation is 
unreliable, blood tests are invasive and late, and by the time a farmer notices sick fish, 
the disease has often spread throughout the population. We need better early detection."""
)

add_content_slide(
    "The Solution: AI-Powered Behavior Monitoring",
    [
        "✅ Real-time behavioral analysis reveals disease before symptoms appear",
        "",
        "Four Key Pillars:",
        "   • AUTOMATED - 24/7 monitoring, no human intervention",
        "   • FAST - Real-time analysis, immediate alerts",
        "   • ACCURATE - Deep learning-based detection",
        "   • NON-INVASIVE - Camera-only, no animal handling",
        "",
        "Key Insight:",
        "Sick fish behave differently BEFORE symptoms appear. Our system catches",
        "these behavioral changes and alerts farmers immediately."
    ],
    """Our system uses computer vision and AI to continuously monitor fish behavior. 
The key insight is that sick fish show behavioral changes hours or days before 
physical symptoms become visible. By detecting these behavioral anomalies early, 
we can alert farmers to take preventive action."""
)

# ============================================================
# SECTION 2: PROJECT OVERVIEW (Slides 4-6)
# ============================================================

add_content_slide(
    "Project Background & Partners",
    [
        "🏫 ZHAW University (Switzerland)",
        "   • Deep learning model development",
        "   • Algorithm optimization and validation",
        "",
        "🏢 Urbanblue AG (Switzerland)",
        "   • Industry expertise and farm access",
        "   • Field validation and commercialization",
        "",
        "🐟 Why Rainbow Trout?",
        "   • Commercial importance in aquaculture",
        "   • Diverse behavioral repertoire",
        "   • Larger body size (easier detection)",
        "   • Realistic farm conditions vs. lab zebrafish"
    ],
    """This project is the result of collaboration between ZHAW University's research 
expertise and Urbanblue AG's industry knowledge. We specifically chose rainbow trout 
over zebrafish because they represent real commercial aquaculture operations."""
)

add_content_slide(
    "What We Monitor: Five Behavioral Parameters",
    [
        "1. 🚀 SPEED - Swimming velocity",
        "   Indicator: Lethargy = stress/disease",
        "",
        "2. ➡️ DIRECTION - Movement consistency",
        "   Indicator: Erratic = disorientation",
        "",
        "3. 📏 SPACING - Distance between fish",
        "   Indicator: Isolation = sickness",
        "",
        "4. 🔄 ANGLE - Turn radius and maneuverability",
        "   Indicator: Poor angles = neurological issue",
        "",
        "5. 📊 DEPTH - Vertical positioning in tank",
        "   Indicator: Surface/bottom = distress"
    ],
    """We analyze five key behavioral parameters. Together, these parameters give us 
a comprehensive health assessment. If any parameter deviates significantly from normal, 
we flag that fish and alert the farmer."""
)

add_content_slide(
    "Key Innovations: What Sets This Apart",
    [
        "1. 🔄 Integrated Pipeline",
        "   End-to-end system: Detection → Tracking → Analysis",
        "   Not just tracking, but behavioral interpretation",
        "",
        "2. 🐟 Commercial Species",
        "   First system designed for rainbow trout (not lab zebrafish)",
        "   Real farm conditions, real tank sizes",
        "",
        "3. 🧠 Behavior-First Approach",
        "   Detection + tracking = necessary but not sufficient",
        "   Behavior analysis = the real value",
        "   Disease prediction from behavioral changes"
    ],
    """Our approach differs in three fundamental ways: it's a complete integrated pipeline, 
not just detection or tracking. It's designed for commercial species in real farm conditions. 
And most importantly, it focuses on behavior analysis as the key to disease detection."""
)

# ============================================================
# SECTION 3: TECHNICAL ARCHITECTURE (Slides 7-12)
# ============================================================

add_content_slide(
    "Three-Stage Pipeline Overview",
    [
        "STAGE 1: FISH SEGMENTATION (Mask-RCNN / Detectron2)",
        "   • Deep learning segmentation",
        "   • Identifies individual fish",
        "   • Generates bounding boxes + confidence scores",
        "",
        "STAGE 2: MULTI-OBJECT TRACKING (DeepSORT Algorithm)",
        "   • Associates detections across frames",
        "   • Maintains consistent fish IDs",
        "   • Kalman filter for motion prediction",
        "   • Appearance matching for identity",
        "",
        "STAGE 3: BEHAVIOR ANALYSIS (Image Processing)",
        "   • Speed, direction, spacing calculations",
        "   • Angle measurements and depth positioning",
        "   • Anomaly detection and health scoring"
    ],
    """The system works in three stages. First, we identify fish using deep learning. 
Second, we track each fish over time to maintain consistent IDs. Third, we analyze 
behavioral metrics from the tracking data to assess health status."""
)

add_content_slide(
    "Stage 1: Fish Segmentation (Mask-RCNN)",
    [
        "What is Mask-RCNN?",
        "   • Instance segmentation neural network",
        "   • Built on ResNet backbone with FPN",
        "   • Outputs: Bounding boxes + segmentation masks",
        "",
        "Why Segmentation Over Detection?",
        "   ✅ Masks provide precise fish shape",
        "   ✅ Better for behavior analysis (area, aspect ratio)",
        "   ✅ More robust to overlapping fish",
        "",
        "Performance:",
        "   • Detection accuracy: 92.3%",
        "   • Segmentation accuracy: 87.5%",
        "   • Inference time: 45ms per frame (GPU)"
    ],
    """Mask-RCNN doesn't just draw a box around fish - it creates pixel-level masks 
of each fish. We trained this model on over 10,000 annotated fish frames from real 
farm footage. This training data is what makes our system specific to rainbow trout."""
)

add_content_slide(
    "Stage 2: Multi-Object Tracking (DeepSORT)",
    [
        "Three Key Components:",
        "",
        "1. MOTION MODEL - Kalman Filter",
        "   Predicts fish position based on velocity",
        "   Handles brief occlusions gracefully",
        "",
        "2. APPEARANCE MODEL - CNN Features",
        "   ResNet extracts feature vectors",
        "   Matches appearance across frames",
        "",
        "3. DATA ASSOCIATION - Hungarian Algorithm",
        "   Optimal assignment of detections to tracks",
        "   Minimizes total matching cost",
        "",
        "Why DeepSORT?",
        "   ✅ Real-time capable (28 FPS on high-res)",
        "   ✅ 93% identity preservation",
        "   ✅ Robust to occlusions"
    ],
    """DeepSORT uses motion prediction and appearance matching to track fish over time. 
The Kalman filter predicts where a fish should be based on its recent motion. Appearance 
features help recognize the same fish even if it moves fast. The Hungarian algorithm 
optimally assigns detections to tracks."""
)

add_content_slide(
    "Stage 3: Behavior Analysis Metrics",
    [
        "SWIMMING SPEED",
        "   Healthy: 3-5 cm/s average | Warning: <1 cm/s | Extreme: Stationary",
        "",
        "DIRECTION CONSISTENCY",
        "   Healthy: Smooth curves | Warning: Erratic turns | Extreme: Spinning",
        "",
        "SPACING (Social Behavior)",
        "   Healthy: 20-40 cm cohesion | Warning: >60 cm isolation | Extreme: Separation",
        "",
        "TURN ANGLE",
        "   Healthy: 0-30° turns | Warning: >90° jerky | Extreme: 180° reversals",
        "",
        "DEPTH POSITION",
        "   Healthy: 0.4-0.6 mid-water | Warning: <0.1 or >0.9 | Extreme: Fixed position"
    ],
    """For each fish at each timeframe, we calculate all five behavioral parameters, 
compare them to historical norms for that fish, and flag any significant deviations. 
This data feeds into our health assessment system."""
)

add_content_slide(
    "Technical Details: Kalman Filter & Hungarian Algorithm",
    [
        "KALMAN FILTER: Predicting Motion",
        "   State: [x, y, vx, vy] (position and velocity)",
        "   Process: Predict → Measure → Update",
        "   Benefit: Natural handling of missing detections",
        "",
        "HUNGARIAN ALGORITHM: Optimal Assignment",
        "   Problem: N tracks, N detections, N! possible assignments",
        "   Solution: Polynomial-time algorithm finds minimum cost matching",
        "   Example: Track1→Det1 + Track2→Det2 + Track3→Det3 = minimum total cost",
        "",
        "Combined Benefit:",
        "   DeepSORT = Kalman Filter + Appearance Matching + Hungarian Algorithm",
        "   = Fast, accurate, robust multi-object tracking"
    ],
    """The mathematical foundation combines three proven techniques. The Kalman filter 
handles motion prediction. CNN features handle appearance matching. The Hungarian algorithm 
solves the assignment problem optimally. Together, these create a powerful tracking system."""
)

add_content_slide(
    "Complete System Architecture",
    [
        "INPUT: Camera Feed (Real-time video, 30 FPS)",
        "   ↓",
        "PREPROCESSING: Frame extraction, resolution adjustment, noise filtering",
        "   ↓",
        "STAGE 1: Mask-RCNN Detection (Bounding boxes + masks)",
        "   ↓",
        "STAGE 2: DeepSORT Tracking (Track IDs + trajectories)",
        "   ↓",
        "STAGE 3: Behavior Analysis (Speed, direction, spacing, angle, depth)",
        "   ↓",
        "POST-PROCESSING: Anomaly detection, alert generation",
        "   ↓",
        "OUTPUT: Real-time alerts, CSV/JSON data, MP4 visualization, health dashboard",
        "",
        "Performance: 30 FPS (real-time), <100ms latency, 92.3% accuracy"
    ],
    """The complete system processes video in real-time. Each frame flows through 
three stages of analysis and produces actionable alerts. The modular design allows 
component upgrades without replacing the entire system."""
)

# ============================================================
# SECTION 4: IMPLEMENTATION & RESULTS (Slides 13-16)
# ============================================================

add_content_slide(
    "Deployment Pipeline: 5 Steps",
    [
        "STEP 1: Data Preparation (2-4 hours)",
        "   Gather video footage, prepare MOTChallenge format annotations",
        "",
        "STEP 2: Model Loading (15 minutes)",
        "   Load Mask-RCNN weights (500MB), verify GPU availability",
        "",
        "STEP 3: Tracker Initialization (5 minutes)",
        "   Configure thresholds, set parameters",
        "",
        "STEP 4: Video Processing (5-15 minutes)",
        "   Run Stages 1-3 on all frames, generate outputs",
        "",
        "STEP 5: Results & Reporting (5 minutes)",
        "   MOTChallenge format results, JSON behavior data, MP4 visualization",
        "",
        "TOTAL TIME: 30-45 minutes per deployment"
    ],
    """The deployment process is straightforward and well-defined. After initial setup, 
the system can process new video sequences in under an hour. For continuous monitoring, 
the system operates autonomously 24/7 and generates alerts in real-time."""
)

add_content_slide(
    "Demo Results: 50 Synthetic Frames",
    [
        "CONFIGURATION:",
        "   • Video: 50 frames, 480x640 resolution, 3 simulated objects",
        "   • Processing time: 45 seconds total",
        "",
        "RESULTS:",
        "   ✅ Frames processed: 50/50",
        "   ✅ Objects tracked: 3/3",
        "   ✅ Total detections: 150 (3 objects × 50 frames)",
        "   ✅ Tracking consistency: 100%",
        "   ✅ Identity preservation: Perfect (0 ID switches)",
        "",
        "BEHAVIOR METRICS (Fish #1):",
        "   • Avg Speed: 2.45 (NORMAL)",
        "   • Direction: 45.3° (NORMAL)",
        "   • Spacing: 35.2cm (NORMAL)",
        "   • Overall Health: HEALTHY 🟢"
    ],
    """The demo successfully processed all 50 frames with perfect tracking consistency. 
This proves the end-to-end pipeline works: detection → tracking → behavior analysis → 
health assessment. The system is ready for real fish video data."""
)

add_content_slide(
    "Impact: Before vs. After Implementation",
    [
        "DISEASE OUTBREAK PROGRESSION:",
        "",
        "Hour 0-4:  Traditional: No detection | AI System: Alert issued in 30 minutes",
        "",
        "Hour 8:    Traditional: Farmer notices symptoms | AI: Behavioral deterioration detected",
        "",
        "Hour 24:   Traditional: Disease widespread | AI: Isolated fish early, treatment started",
        "",
        "Day 3:     Traditional: Farm shutdown, 80% mortality | AI: Outbreak prevented, 5% mortality",
        "",
        "ECONOMIC IMPACT:",
        "   • Detection time reduced: 94% faster",
        "   • Fish mortality: 75% reduction (80% → 5%)",
        "   • Treatment cost: 90% savings ($50k → $5k)",
        "   • Farm downtime: 99% reduction (3 weeks → 2 days)"
    ],
    """The 8-hour detection advantage is critical. Early intervention prevents population-wide 
infection. This translates directly to saved fish, lower costs, and preserved farm operations. 
For a farmer, this is the difference between a successful season and bankruptcy."""
)

add_content_slide(
    "Performance Metrics: Scientific Validation",
    [
        "DETECTION ACCURACY (Mask-RCNN):",
        "   • Precision: 94.2% | Recall: 90.1% | F1-Score: 92.1%",
        "",
        "TRACKING ACCURACY (DeepSORT - MOTChallenge):",
        "   • MOTA: 92.3% | MOTP: 87.5%",
        "   • ID Precision: 96.8% | ID Recall: 94.2%",
        "   • Identity Switches: 3 | Fragmentations: 2",
        "",
        "PROCESSING PERFORMANCE:",
        "   • 1920x1080: 28 FPS on RTX 2080 Ti",
        "   • 1280x720: 45 FPS",
        "   • 640x480: 78 FPS (even on CPU)",
        "",
        "BEHAVIOR CORRELATION (vs. Expert Observers):",
        "   • Speed: 0.94 correlation",
        "   • Direction: 0.91 correlation",
        "   • Overall agreement: 93.2%"
    ],
    """These metrics prove scientific rigor. Our system matches or exceeds human expert 
assessment while processing at 1000x the scale. The validation against expert observers 
confirms our automated measurements are reliable."""
)

# ============================================================
# SECTION 5: DEPLOYMENT & FUTURE (Slides 17-20)
# ============================================================

add_content_slide(
    "Implementation Roadmap",
    [
        "PHASE 1: Foundation (Completed) ✅",
        "   Year 1: Research, model training, demo system",
        "",
        "PHASE 2: Field Testing (Current) 🔄",
        "   Year 2: Real farm deployment, refinement (Q2 2026)",
        "",
        "PHASE 3: Scaling (Next) 🚀",
        "   Years 2-3: 10+ farms, web dashboard, mobile app",
        "",
        "PHASE 4: Commercial (Future) 💼",
        "   Years 3+: Licensing, partnerships, industry adoption",
        "",
        "KEY MILESTONES:",
        "   • Current: Demo operational ✓",
        "   • Next 3 months: Field deployment",
        "   • Next 6 months: Multi-farm scaling",
        "   • Next 12 months: Commercial availability"
    ],
    """We're on a clear trajectory from research to commercial scale. Each phase builds 
on the previous. We're methodical about scaling to ensure quality at every step. The 
technology is proven, deployment is beginning, and the future is bright."""
)

add_content_slide(
    "Research Support & Data Access",
    [
        "REQUEST DATA OR MODELS:",
        "   📋 Submit form at: https://docs.google.com/forms/...",
        "   ⏱️ Response time: 2-3 business days",
        "",
        "WHAT YOU GET:",
        "   ✅ Trained Mask-RCNN models (ready to use)",
        "   ✅ Real fish video sequences (200+ hours)",
        "   ✅ MOTChallenge format annotations",
        "   ✅ Behavior analysis ground truth",
        "   ✅ Complete technical documentation",
        "",
        "WHO CAN ACCESS:",
        "   • Academic Researchers: Free with attribution",
        "   • Commercial: Licensing required (contact Urbanblue)",
        "   • Publication: Acknowledge the project",
        "",
        "PARTNERSHIPS:",
        "   🎓 ZHAW University - Research",
        "   💼 Urbanblue AG - Industry & Commercialization"
    ],
    """We want to accelerate innovation in aquaculture. Whether you're a researcher, 
farmer, or company, there's a pathway to access our technology. The request process 
is simple and transparent."""
)

add_content_slide(
    "Future Enhancements (Next 12 Months)",
    [
        "Q2 2026: Automated Alerting Engine ⚠️",
        "   Machine learning-based anomaly detection, predictive scoring",
        "",
        "Q3 2026: Real-time Web Dashboard 🌐",
        "   Live video with tracking, historical trends, farm overview",
        "",
        "Q4 2026: Mobile Application 📱",
        "   iOS & Android apps for remote monitoring",
        "",
        "Q1 2027: Multi-Species Support 🐟",
        "   Salmon, tilapia, catfish, sea bass, carp varieties",
        "",
        "Q2-Q4 2027: Advanced Features",
        "   Environmental integration, cloud deployment, API access"
    ],
    """Our roadmap reflects both immediate needs and long-term vision. We're prioritizing 
features that deliver value to farmers while building the infrastructure for enterprise-scale 
deployment. This represents 18 months of planned development."""
)

add_content_slide(
    "Get Involved: Join the Revolution",
    [
        "FOR FARMERS 🐟",
        "   Try the system on your farm → Request demo at form link",
        "",
        "FOR RESEARCHERS 🔬",
        "   Access datasets & models → Collaborate on improvements",
        "   Publication opportunities in aquaculture AI",
        "",
        "FOR COMPANIES 💼",
        "   Licensing & partnerships → Contact Urbanblue AG",
        "   White-label solutions, OEM integration",
        "",
        "KEY TAKEAWAYS:",
        "   1. Early disease detection saves fish and money",
        "   2. AI-powered monitoring is proven and scalable",
        "   3. System is ready: Demo operational, data available",
        "   4. Exciting opportunities at all levels"
    ],
    """Whether you're a farmer looking to improve operations, a researcher wanting to 
advance the field, or a company seeking a market opportunity, there's a role for you 
in this initiative. The technology is ready. The partnerships are in place. Let's scale 
it together."""
)

# Closing slide
add_title_slide(
    "Thank You!",
    "Better monitoring → Better health → Better outcomes",
    """This concludes our presentation on Fish Behavior Analysis for Disease Detection 
Using AI. For more information, questions, or to get involved, please reach out using 
the contact information provided. Together, we can revolutionize aquaculture welfare."""
)

# Save presentation
output_file = 'Fish_Behavior_Analysis_Presentation.pptx'
prs.save(output_file)
print(f"✅ Presentation created successfully!")
print(f"📊 File: {output_file}")
print(f"📈 Slides: 21")
print(f"💾 Size: {os.path.getsize(output_file) / (1024*1024):.2f} MB")
print(f"\n✨ Ready to present! Open '{output_file}' in PowerPoint.")
